import collections.abc  # Import to avoid potential errors with collections
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE  # Shape type constants
from pptx.util import Cm
import MeCab
import pandas as pd

# プレゼンテーションファイルを読み込み
ppt = Presentation("./sample2.pptx")
wakati = MeCab.Tagger("-Owakati")
result = []

# スライドごとにテキストと画像の情報を取得
for slide_num, slide in enumerate(ppt.slides, start=1):
    print(f"Slide {slide_num}")
    
    for shape in slide.shapes:
        # テキストがある場合の処理
        if shape.has_text_frame:
            left, top = shape.left, shape.top
            for paragraph in shape.text_frame.paragraphs:
                print(f"Text: {paragraph.text}")
                print(f"Position - Left: {left / Cm(1):.2f} cm, Top: {top / Cm(1):.2f} cm")
                result.append(wakati.parse(paragraph.text).split())

                                
                for run in paragraph.runs:
                    try:
                        # フォントカラーの取得
                        font_color = run.font.color
                        color_rgb = font_color.rgb if font_color and font_color.rgb else "デフォルト色"
                        print(f"Color: {color_rgb}")
                    except AttributeError:
                        print("Color: None")
                    
                    # フォントサイズの取得
                    font_size = run.font.size
                    font_size_pt = font_size.pt if font_size else "None"
                    print(f"Font Size: {font_size_pt} pt")

        # 図や写真の場合の処理
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            left, top = shape.left, shape.top
            width, height = shape.width, shape.height
            image = shape.image
            print("Image:")
            print(f"  File name: {image.filename}")
            print(f"  Position - Left: {left / Cm(1):.2f} cm, Top: {top / Cm(1):.2f} cm")
            print(f"  Size - Width: {width / Cm(1):.2f} cm, Height: {height / Cm(1):.2f} cm")
            
            # 画像ファイルを保存する場合
            # with open(f"extracted_image_{slide_num}.jpg", "wb") as img_file:
            #     img_file.write(image.blob)
    
    print("\n")  # スライドの区切り
df = pd.DataFrame(result)    
df.to_csv('output.csv', index=False, encoding='utf-8')
print(df)