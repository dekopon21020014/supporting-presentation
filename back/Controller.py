import os
from fastapi import APIRouter, File, UploadFile, status
from fastapi.responses import ORJSONResponse
from pptx import Presentation
from io import BytesIO
import MeCab
from collections import Counter

router = APIRouter(tags=["Demo"], default_response_class=ORJSONResponse)


@router.get(
    path="/",
    status_code=status.HTTP_200_OK,
    responses={
        status.HTTP_200_OK: {
            "description": "200 OK",
            "content": {
                "application/json": {
                    "example": {"message": "yahoo.co.jp"}
                }
            }
        },
        status.HTTP_500_INTERNAL_SERVER_ERROR: {
            "description": "Internal server error dayo"
        },
    },
)
async def get_demo():
    return {"message": "hello"}


@router.post(
    path="/upload",
    status_code=status.HTTP_200_OK,
    responses={
        status.HTTP_200_OK: {
            "description": "200 OK",
        },
        status.HTTP_500_INTERNAL_SERVER_ERROR: {
            "description": "Internal server error"
        },
    },
)
async def upload_pptx(file: UploadFile = File(...)):
    # ファイルの読み込み
    file_content = await file.read()

    # プレゼンテーションを読み込み
    presentation = Presentation(BytesIO(file_content))
    all_text = " ".join(
        shape.text.strip()
        for slide in presentation.slides
        for shape in slide.shapes
        if hasattr(shape, "text") and shape.text.strip()
    )

    # 形態素解析と品詞カウント
    pos_counts = Counter()  # 品詞のカウント
    surface_counts = {}  # 単語のごとのカウント
    '''
    surface_countsはこんな感じになる
    surface_counts = {
        'BOS/EOS': Counter({'': 2}), 
        '名詞': Counter({'30': 9, '回': 7, '発表': 6}),
        '接尾辞': Counter({'日': 6, 'さん': 4}),
        '助詞': Counter({'の': 23, 'て': 14}),
        '接頭辞': Counter({'第': 5, 'お': 3, 'ご': 1}), 
        '補助記号': Counter({'、': 16, '（': 9, '）': 9}),
        '助動詞': Counter({'た': 4, 'ます': 3}),
        '記号': Counter({'-': 6, 'ち': 1}), 
        '動詞': Counter({'し': 7, 'いる': 4}),
        '副詞': Counter({'ぜひ': 2, 'もう': 1}), 
        '形状詞': Counter({'たくさん': 1, '可能': 1}), 
        '代名詞': Counter({'どこ': 1, 'ここ': 1}), 
        '連体詞': Counter({'この': 1, 'その': 1}), 
        '感動詞': Counter({'ん': 1}), 
        '形容詞': Counter({'小さく': 1})}
    }
    '''

    mecab = MeCab.Tagger()
    node = mecab.parseToNode(all_text)
    while node:
        surface = node.surface
        pos = node.feature.split(",")[0]
        pos_counts[pos] += 1

        if pos not in surface_counts:
            surface_counts[pos] = Counter()
        surface_counts[pos][surface] += 1
        node = node.next

    # 上位3位の名詞、動詞、助詞を取得
    top_nouns = surface_counts['名詞'].most_common(
        3)  # noun_counts.most_common(3)
    top_verbs = surface_counts['動詞'].most_common(
        3)  # verb_counts.most_common(3)
    top_particles = surface_counts['助詞'].most_common(
        3)  # particle_counts.most_common(3)
    print(surface_counts)
    # 結果を返す
    return {
        "message": "File processed successfully",
        "slide_count": len(presentation.slides),
        "pos_counts": pos_counts,  # 全品詞のカウント
        "top_nouns": top_nouns,  # 名詞のトップ3
        "top_verbs": top_verbs,  # 動詞のトップ3
        "top_particles": top_particles,  # 助詞のトップ3
    }
